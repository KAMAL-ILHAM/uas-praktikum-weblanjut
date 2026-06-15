from PyPDF2 import PdfReader
import docx
from pptx import Presentation

def ekstrak_teks_dari_file(file_obj):
    """
    Fungsi untuk membaca teks dari file PDF, DOCX, atau PPTX
    yang diunggah oleh user (InMemoryUploadedFile).
    """
    teks_lengkap = ""
    nama_file = file_obj.name.lower()

    try:
        # 1. JIKA FILE PDF
        if nama_file.endswith('.pdf'):
            reader = PdfReader(file_obj)
            for page in reader.pages:
                teks_terbaca = page.extract_text()
                if teks_terbaca:
                    teks_lengkap += teks_terbaca + "\n"
                    
        # 2. JIKA FILE WORD (.docx)
        elif nama_file.endswith('.docx'):
            doc = docx.Document(file_obj)
            for para in doc.paragraphs:
                teks_lengkap += para.text + "\n"
                
        # 3. JIKA FILE POWERPOINT (.pptx)
        elif nama_file.endswith('.pptx'):
            prs = Presentation(file_obj)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        teks_lengkap += shape.text + "\n"
        
        # 4. JIKA FORMAT LAIN
        else:
            raise ValueError("Format file tidak didukung. Harap gunakan PDF, DOCX, atau PPTX.")
            
    except Exception as e:
        raise Exception(f"Terjadi kesalahan saat mengekstrak teks: {str(e)}")
        
    return teks_lengkap.strip()